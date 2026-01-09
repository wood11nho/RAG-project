import os
import shutil
import subprocess
import pickle
import logging
from pathlib import Path
from typing import List, Optional, Generator

from dotenv import load_dotenv
from pptx import Presentation

# LangChain & Vector Store
from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document
from langchain_community.retrievers import BM25Retriever

# Configure Logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    datefmt="%H:%M:%S"
)
logger = logging.getLogger(__name__)

# Load Environment Variables
load_dotenv()

# --- Configuration ---
# Using pathlib for cross-platform compatibility
BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
DB_DIR = BASE_DIR / "chroma_db"
TEMP_DIR = BASE_DIR / "temp_conversion"
BM25_PATH = DB_DIR / "bm25_retriever.pkl"

COLLECTION_NAME = "course_materials"
EMBEDDING_MODEL = "all-MiniLM-L6-v2"


class DocumentProcessor:
    """
    Handles the ingestion of course materials (PPTX) into vector and sparse indices.
    """

    def __init__(self, source_dir: Path, persist_dir: Path):
        self.source_dir = source_dir
        self.persist_dir = persist_dir
        
        logger.info(f"Loading embedding model: {EMBEDDING_MODEL}")
        self.embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

        # Ensure temp directory exists for legacy conversions
        if not TEMP_DIR.exists():
            TEMP_DIR.mkdir(parents=True)

    def _convert_ppt_to_pptx(self, ppt_path: Path) -> Optional[Path]:
        """
        Converts legacy .ppt files to .pptx using LibreOffice.
        Returns the path to the converted file or None if failed.
        """
        output_filename = ppt_path.name.replace(".ppt", ".pptx")
        output_path = TEMP_DIR / output_filename

        if output_path.exists():
            logger.debug(f"Cached conversion found: {output_path}")
            return output_path

        logger.info(f"Converting legacy format: {ppt_path.name}")
        
        # Headless conversion via LibreOffice
        cmd = [
            "libreoffice", "--headless", "--convert-to", "pptx",
            "--outdir", str(TEMP_DIR), str(ppt_path)
        ]
        
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, check=False)
            if result.returncode == 0 and output_path.exists():
                return output_path
            
            logger.error(f"Conversion failed for {ppt_path.name}: {result.stderr}")
            return None
        except FileNotFoundError:
            logger.warning("LibreOffice not found. Skipping .ppt file.")
            return None

    def _clean_text(self, text: str) -> str:
        """Normalizes whitespace in extracted text."""
        return " ".join(text.split()) if text else ""

    def _extract_slide_content(self, slide, slide_number: int) -> str:
        """Extracts text from title, body, and notes of a single slide."""
        content_parts = []

        # 1. Title
        if slide.shapes.title and slide.shapes.title.text:
            title_text = self._clean_text(slide.shapes.title.text)
            content_parts.append(f"Title: {title_text}")

        # 2. Body Text
        body_text = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if hasattr(shape, "text") and shape.text.strip():
                body_text.append(self._clean_text(shape.text))
        
        if body_text:
            content_parts.append("Content: " + " | ".join(body_text))

        # 3. Speaker Notes (Crucial for context)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                content_parts.append(f"Notes: {self._clean_text(notes)}")

        return "\n".join(content_parts)

    def process_file(self, file_path: Path) -> List[Document]:
        """Parses a single PPTX file into LangChain Documents."""
        docs = []
        try:
            prs = Presentation(file_path)
            logger.info(f"Processing: {file_path.name} ({len(prs.slides)} slides)")
            
            for i, slide in enumerate(prs.slides):
                text = self._extract_slide_content(slide, i + 1)
                
                # Skip empty or near-empty slides
                if len(text) < 15:
                    continue

                metadata = {
                    "source": file_path.name,
                    "slide_number": i + 1,
                    "file_type": "pptx"
                }
                docs.append(Document(page_content=text, metadata=metadata))
                
        except Exception as e:
            logger.error(f"Failed to process {file_path.name}: {e}")
            
        return docs

    def run(self):
        """Main execution pipeline."""
        if not self.source_dir.exists():
            logger.error(f"Source directory '{self.source_dir}' does not exist.")
            return

        all_documents = []

        # 1. Handle Legacy .ppt files
        ppt_files = list(self.source_dir.glob("*.ppt"))
        if ppt_files:
            logger.info(f"Found {len(ppt_files)} legacy .ppt files.")
            for ppt in ppt_files:
                converted_path = self._convert_ppt_to_pptx(ppt)
                if converted_path:
                    all_documents.extend(self.process_file(converted_path))

        # 2. Handle Modern .pptx files
        pptx_files = list(self.source_dir.glob("*.pptx"))
        for pptx in pptx_files:
            all_documents.extend(self.process_file(pptx))

        if not all_documents:
            logger.warning("No documents extracted. Exiting.")
            return

        logger.info(f"Total extracted documents: {len(all_documents)}")

        # 3. Vector Store Ingestion (Chroma)
        logger.info(f"Persisting Vector Store to: {self.persist_dir}")
        Chroma.from_documents(
            documents=all_documents,
            embedding=self.embeddings,
            persist_directory=str(self.persist_dir),
            collection_name=COLLECTION_NAME
        )

        # 4. Sparse Index Ingestion (BM25)
        logger.info("Building and saving BM25 index...")
        bm25_retriever = BM25Retriever.from_documents(all_documents)
        
        with open(BM25_PATH, "wb") as f:
            pickle.dump(bm25_retriever, f)

        # Cleanupconda env
        if TEMP_DIR.exists():
            shutil.rmtree(TEMP_DIR)
            
        logger.info("Ingestion pipeline completed successfully.")


if __name__ == "__main__":
    processor = DocumentProcessor(DATA_DIR, DB_DIR)
    processor.run()